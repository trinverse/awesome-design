#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

def create_network_architecture_presentation():
    """Create PowerPoint presentation with network architecture diagrams"""
    prs = Presentation()
    
    # Set slide size to 16:9 wide
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HAL Network Architecture"
    subtitle.text = "Logical and Physical Network Diagrams\nIntegrated Asset Management & Monitoring Platform"
    
    # Create Logical Network Architecture slide
    create_logical_network_diagram(prs)
    
    # Create Physical Network Architecture slide
    create_physical_network_diagram(prs)
    
    # Save presentation
    prs.save('HAL_Network_Architecture_Diagrams.pptx')
    print("PowerPoint presentation created successfully: HAL_Network_Architecture_Diagrams.pptx")

def create_logical_network_diagram(prs):
    """Create logical network architecture diagram"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "HAL Logical Network Architecture"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Define colors for different zones
    colors = {
        'dmz': RGBColor(255, 200, 200),        # Light red
        'web_app': RGBColor(200, 255, 200),    # Light green
        'mz_db': RGBColor(200, 200, 255),      # Light blue
        'mgmt': RGBColor(255, 255, 200),       # Light yellow
        'external': RGBColor(240, 240, 240),   # Light gray
        'firewall': RGBColor(255, 150, 100),   # Orange
        'router': RGBColor(100, 150, 255),     # Blue
        'lb': RGBColor(150, 255, 150),         # Green
    }
    
    # External zone
    ext_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1), Inches(12.3), Inches(1.5))
    ext_zone.fill.solid()
    ext_zone.fill.fore_color.rgb = colors['external']
    ext_zone.line.color.rgb = RGBColor(100, 100, 100)
    add_text_to_shape(ext_zone, "External Network", Pt(12), bold=True)
    
    # Internet and MPLS routers
    internet_router1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.3), Inches(1.5), Inches(0.8))
    internet_router1.fill.solid()
    internet_router1.fill.fore_color.rgb = colors['router']
    add_text_to_shape(internet_router1, "Internet\nRouter 1\nISR 4461", Pt(10))
    
    internet_router2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(1.3), Inches(1.5), Inches(0.8))
    internet_router2.fill.solid()
    internet_router2.fill.fore_color.rgb = colors['router']
    add_text_to_shape(internet_router2, "Internet\nRouter 2\nISR 4461", Pt(10))
    
    mpls_router1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.3), Inches(1.5), Inches(0.8))
    mpls_router1.fill.solid()
    mpls_router1.fill.fore_color.rgb = colors['router']
    add_text_to_shape(mpls_router1, "MPLS\nRouter 1\nISR 4461", Pt(10))
    
    mpls_router2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(1.3), Inches(1.5), Inches(0.8))
    mpls_router2.fill.solid()
    mpls_router2.fill.fore_color.rgb = colors['router']
    add_text_to_shape(mpls_router2, "MPLS\nRouter 2\nISR 4461", Pt(10))
    
    # Link Load Balancer
    llb = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(5.5), Inches(2.8), Inches(2), Inches(0.8))
    llb.fill.solid()
    llb.fill.fore_color.rgb = colors['lb']
    add_text_to_shape(llb, "Link Load\nBalancer", Pt(11), bold=True)
    
    # Perimeter Firewall
    perimeter_fw = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(5), Inches(3.8), Inches(3), Inches(0.8))
    perimeter_fw.fill.solid()
    perimeter_fw.fill.fore_color.rgb = colors['firewall']
    add_text_to_shape(perimeter_fw, "Perimeter Firewall\nPA-5410", Pt(11), bold=True)
    
    # DMZ Zone
    dmz_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.8), Inches(3), Inches(1.5))
    dmz_zone.fill.solid()
    dmz_zone.fill.fore_color.rgb = colors['dmz']
    dmz_zone.line.color.rgb = RGBColor(200, 100, 100)
    add_text_to_shape(dmz_zone, "DMZ Zone (VLAN 103)", Pt(11), bold=True)
    
    # Sterling File Gateway in DMZ
    sfg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.3), Inches(2), Inches(0.7))
    sfg.fill.solid()
    sfg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_text_to_shape(sfg, "Sterling File\nGateway (SFG)", Pt(10))
    
    # Core Firewall
    core_fw = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(5), Inches(5.2), Inches(3), Inches(0.8))
    core_fw.fill.solid()
    core_fw.fill.fore_color.rgb = colors['firewall']
    add_text_to_shape(core_fw, "Core Firewall\nFG-400F", Pt(11), bold=True)
    
    # Web/App Zone
    web_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.8), Inches(4.5), Inches(1.5))
    web_zone.fill.solid()
    web_zone.fill.fore_color.rgb = colors['web_app']
    web_zone.line.color.rgb = RGBColor(100, 200, 100)
    add_text_to_shape(web_zone, "Web/App Zone", Pt(11), bold=True)
    
    # Server Load Balancer
    slb = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(5.5), Inches(5.3), Inches(1.5), Inches(0.7))
    slb.fill.solid()
    slb.fill.fore_color.rgb = colors['lb']
    add_text_to_shape(slb, "Server LB", Pt(10))
    
    # Management Zone
    mgmt_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9), Inches(4.8), Inches(3.3), Inches(1.5))
    mgmt_zone.fill.solid()
    mgmt_zone.fill.fore_color.rgb = colors['mgmt']
    mgmt_zone.line.color.rgb = RGBColor(200, 200, 100)
    add_text_to_shape(mgmt_zone, "Management Zone (VLAN 104)", Pt(11), bold=True)
    
    # Management components
    mgmt_components = [
        ("vPAS NMS", Inches(9.2), Inches(5.3)),
        ("SevOne", Inches(10.3), Inches(5.3)),
        ("vDNC", Inches(11.4), Inches(5.3))
    ]
    
    for comp_name, left, top in mgmt_components:
        comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.9), Inches(0.6))
        comp.fill.solid()
        comp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        add_text_to_shape(comp, comp_name, Pt(9))
    
    # MZ/DB Zone
    db_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.5), Inches(11.8), Inches(0.8))
    db_zone.fill.solid()
    db_zone.fill.fore_color.rgb = colors['mz_db']
    db_zone.line.color.rgb = RGBColor(100, 100, 200)
    add_text_to_shape(db_zone, "MZ/DB Zone (VLAN 105) - Database & Backend Services", Pt(11), bold=True)
    
    # Add application components in Web/App zone
    app_y = 5.4
    app_components = [
        ("Maximo/MAS", Inches(4.2)),
        ("Instana", Inches(5.2)),
        ("Kafka", Inches(6.2)),
        ("Reports", Inches(7.2))
    ]
    
    for comp_name, left in app_components:
        comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(app_y), Inches(0.9), Inches(0.5))
        comp.fill.solid()
        comp.fill.fore_color.rgb = RGBColor(255, 255, 255)
        add_text_to_shape(comp, comp_name, Pt(9))
    
    # Add VLAN legend
    legend_left = Inches(0.5)
    legend_top = Inches(7.0)
    legend_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left, legend_top, Inches(6), Inches(0.4))
    legend_box.fill.solid()
    legend_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    legend_text = "VLANs: 101-Network MGMT | 102-Compute MGMT | 103-DMZ | 104-MZ | 105-DB | 106-Security MGMT | 107/108-Internet | 109/110-MPLS | 111-Common Service"
    add_text_to_shape(legend_box, legend_text, Pt(8))
    
    # Add security components legend
    sec_legend_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(7.0), Inches(5.3), Inches(0.4))
    sec_legend_box.fill.solid()
    sec_legend_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    sec_legend_text = "Security: IPS (Intrusion Prevention) | WAF (Web Application Firewall) | SSO with LDAP"
    add_text_to_shape(sec_legend_box, sec_legend_text, Pt(8))

def create_physical_network_diagram(prs):
    """Create physical network architecture diagram"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "HAL Physical Network Architecture - Spine-Leaf Topology"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    
    # Colors for components
    colors = {
        'spine': RGBColor(100, 150, 255),      # Blue
        'leaf': RGBColor(150, 255, 150),       # Green
        'server': RGBColor(200, 200, 200),     # Gray
        'storage': RGBColor(255, 200, 150),    # Orange
        'router': RGBColor(150, 150, 255),     # Light blue
        'firewall': RGBColor(255, 150, 100),   # Orange
    }
    
    # Spine switches at the top
    spine_y = Inches(1.5)
    spine1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), spine_y, Inches(2), Inches(0.8))
    spine1.fill.solid()
    spine1.fill.fore_color.rgb = colors['spine']
    add_text_to_shape(spine1, "Spine Switch 1\nN9K-C9316D-GX", Pt(10), bold=True)
    
    spine2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), spine_y, Inches(2), Inches(0.8))
    spine2.fill.solid()
    spine2.fill.fore_color.rgb = colors['spine']
    add_text_to_shape(spine2, "Spine Switch 2\nN9K-C9316D-GX", Pt(10), bold=True)
    
    # Leaf switches
    leaf_y = Inches(3.5)
    leaf_specs = [
        ("Border Leaf 1\nN9K-C9316D-GX", Inches(1)),
        ("Border Leaf 2\nN9K-C9316D-GX", Inches(3)),
        ("Compute Leaf 1\nN9K-C93180YC-EX", Inches(5)),
        ("Compute Leaf 2\nN9K-C93180YC-EX", Inches(7)),
        ("Service Leaf 1", Inches(9)),
        ("DMZ Leaf 1", Inches(11))
    ]
    
    leaf_shapes = []
    for leaf_name, left in leaf_specs:
        leaf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, leaf_y, Inches(1.8), Inches(0.8))
        leaf.fill.solid()
        leaf.fill.fore_color.rgb = colors['leaf']
        add_text_to_shape(leaf, leaf_name, Pt(9))
        leaf_shapes.append(leaf)
    
    # Add connection lines between spine and leaf (100G)
    add_connection_label(slide, Inches(5), Inches(2.5), "100G Fiber", Pt(9))
    
    # External connectivity
    ext_y = Inches(5)
    
    # Internet routers
    inet_router1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), ext_y, Inches(1.5), Inches(0.7))
    inet_router1.fill.solid()
    inet_router1.fill.fore_color.rgb = colors['router']
    add_text_to_shape(inet_router1, "Internet\nISR 4461", Pt(9))
    
    inet_router2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), ext_y, Inches(1.5), Inches(0.7))
    inet_router2.fill.solid()
    inet_router2.fill.fore_color.rgb = colors['router']
    add_text_to_shape(inet_router2, "Internet\nISR 4461", Pt(9))
    
    # MPLS routers
    mpls_router1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(1.5), Inches(0.7))
    mpls_router1.fill.solid()
    mpls_router1.fill.fore_color.rgb = colors['router']
    add_text_to_shape(mpls_router1, "MPLS\nISR 4461", Pt(9))
    
    mpls_router2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.2), Inches(5.9), Inches(1.5), Inches(0.7))
    mpls_router2.fill.solid()
    mpls_router2.fill.fore_color.rgb = colors['router']
    add_text_to_shape(mpls_router2, "MPLS\nISR 4461", Pt(9))
    
    # Compute servers
    server_y = Inches(5)
    compute_servers = [
        ("OpenShift\nCompute\n25G NICs", Inches(4.5)),
        ("OpenShift\nCompute\n25G NICs", Inches(6.0)),
        ("OpenShift\nCompute\n25G NICs", Inches(7.5))
    ]
    
    for server_name, left in compute_servers:
        server = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, server_y, Inches(1.3), Inches(0.9))
        server.fill.solid()
        server.fill.fore_color.rgb = colors['server']
        server.line.width = Pt(2)
        add_text_to_shape(server, server_name, Pt(8))
    
    add_connection_label(slide, Inches(5.5), Inches(4.5), "25G Fiber", Pt(9))
    
    # Service components
    service_y = Inches(5)
    services = [
        ("Mgmt\nServices\n1G/10G", Inches(9)),
        ("Security\nAppliances\n10G", Inches(10.5))
    ]
    
    for service_name, left in services:
        service = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, service_y, Inches(1.3), Inches(0.9))
        service.fill.solid()
        service.fill.fore_color.rgb = colors['server']
        add_text_to_shape(service, service_name, Pt(8))
    
    # Storage
    storage_y = Inches(6.2)
    storage_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), storage_y, Inches(4.5), Inches(0.8))
    storage_box.fill.solid()
    storage_box.fill.fore_color.rgb = colors['storage']
    add_text_to_shape(storage_box, "SAN Storage - 16G FC Connectivity", Pt(10), bold=True)
    
    # Rack layout information
    rack_info_y = Inches(6.5)
    dc_rack_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), rack_info_y, Inches(5.5), Inches(0.8))
    dc_rack_box.fill.solid()
    dc_rack_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    add_text_to_shape(dc_rack_box, "DC Site: 4 Racks (Network, Compute, Storage, Management)", Pt(10))
    
    dr_rack_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), rack_info_y, Inches(5.5), Inches(0.8))
    dr_rack_box.fill.solid()
    dr_rack_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
    add_text_to_shape(dr_rack_box, "DR Site: 3 Racks (Identical Architecture, Active-Active)", Pt(10))
    
    # Connection speeds legend
    legend_y = Inches(7.0)
    legend_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), legend_y, Inches(11.5), Inches(0.4))
    legend_box.fill.solid()
    legend_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    legend_text = "Connection Speeds: Spine-Leaf: 100G | Compute-Leaf: 25G | Storage: 16G FC | Services: 10G/1G | Dual connections for redundancy"
    add_text_to_shape(legend_box, legend_text, Pt(9))

def add_text_to_shape(shape, text, font_size, bold=False):
    """Add text to a shape with formatting"""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = font_size
        paragraph.font.bold = bold

def add_connection_label(slide, left, top, text, font_size):
    """Add a connection label to the slide"""
    label_box = slide.shapes.add_textbox(left, top, Inches(1), Inches(0.3))
    label_frame = label_box.text_frame
    label_frame.text = text
    label_frame.paragraphs[0].font.size = font_size
    label_frame.paragraphs[0].font.italic = True
    label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    create_network_architecture_presentation()