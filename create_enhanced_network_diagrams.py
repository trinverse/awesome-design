#!/usr/bin/env python3

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
# from pptx.enum.line import MSO_LINE_DASH_STYLE

def create_enhanced_network_presentation():
    """Create enhanced PowerPoint presentation with detailed network architecture diagrams"""
    prs = Presentation()
    
    # Set slide size to 16:9 wide
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HAL Network Architecture"
    subtitle.text = "Logical and Physical Network Diagrams\nIntegrated Asset Management & Monitoring Platform\nSupporting 6 Customer Locations with DC/DR Sites"
    
    # Create overview slide
    create_overview_slide(prs)
    
    # Create Logical Network Architecture slide
    create_enhanced_logical_diagram(prs)
    
    # Create Physical Network Architecture slide
    create_enhanced_physical_diagram(prs)
    
    # Create VLAN and IP scheme slide
    create_vlan_ip_slide(prs)
    
    # Create security architecture slide
    create_security_architecture_slide(prs)
    
    # Create rack diagrams slide
    create_rack_diagrams_slide(prs)
    
    # Save presentation
    prs.save('HAL_Network_Architecture_Enhanced.pptx')
    print("Enhanced PowerPoint presentation created successfully: HAL_Network_Architecture_Enhanced.pptx")

def create_overview_slide(prs):
    """Create overview slide with key architecture highlights"""
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "HAL Network Architecture Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Architecture Features"
    
    # Add bullet points
    bullets = [
        "Spine-Leaf architecture for scalability and low latency",
        "Multi-zone security design with DMZ, Web/App, MZ/DB, and Management zones",
        "Dual data centers (DC & DR) in Active-Active configuration",
        "Support for 6 customer locations via MPLS and Internet connectivity",
        "11 VLANs for network segmentation and security",
        "High-speed connectivity: 100G spine-leaf, 25G compute, 16G FC storage",
        "Comprehensive security stack: Perimeter/Core firewalls, IPS, WAF",
        "Integration with Maximo, SevOne, Instana, and Sterling File Gateway"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1
        p.font.size = Pt(14)

def create_enhanced_logical_diagram(prs):
    """Create enhanced logical network architecture diagram"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "HAL Logical Network Architecture - Multi-Zone Security Design"
    title_frame.paragraphs[0].font.size = Pt(22)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Define enhanced colors
    colors = {
        'internet': RGBColor(240, 240, 255),    # Very light blue
        'dmz': RGBColor(255, 200, 200),         # Light red
        'web_app': RGBColor(200, 255, 200),     # Light green
        'mz_db': RGBColor(200, 220, 255),       # Light blue
        'mgmt': RGBColor(255, 255, 200),        # Light yellow
        'firewall': RGBColor(255, 120, 60),     # Orange-red
        'router': RGBColor(80, 130, 255),       # Blue
        'lb': RGBColor(100, 200, 100),          # Green
        'switch': RGBColor(150, 150, 255),      # Light purple
        'server': RGBColor(220, 220, 220),      # Light gray
    }
    
    # Internet/MPLS cloud at top
    cloud_shape = slide.shapes.add_shape(MSO_SHAPE.CLOUD, Inches(4.5), Inches(0.6), Inches(4.5), Inches(0.8))
    cloud_shape.fill.solid()
    cloud_shape.fill.fore_color.rgb = colors['internet']
    cloud_shape.line.color.rgb = RGBColor(100, 100, 150)
    cloud_shape.line.width = Pt(2)
    add_text_to_shape(cloud_shape, "Internet / MPLS\n6 Customer Locations", Pt(11), bold=True)
    
    # Customer locations
    locations_y = Inches(0.65)
    locations = ["Porbandar", "Ratnagiri", "Chennai", "Bhubaneswar", "Sri Vijaya Puram", "Kochi"]
    for i, loc in enumerate(locations):
        loc_x = Inches(1.5 + i * 1.8)
        loc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, loc_x, locations_y, Inches(1.5), Inches(0.3))
        loc_box.fill.solid()
        loc_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        loc_box.line.color.rgb = RGBColor(150, 150, 150)
        add_text_to_shape(loc_box, loc, Pt(8))
    
    # Border routers
    router_y = Inches(1.7)
    routers = [
        ("Internet Router 1\nISR 4461\nVLAN 107", Inches(2.5)),
        ("Internet Router 2\nISR 4461\nVLAN 108", Inches(4.2)),
        ("MPLS Router 1\nISR 4461\nVLAN 109", Inches(7.5)),
        ("MPLS Router 2\nISR 4461\nVLAN 110", Inches(9.2))
    ]
    
    for router_text, x_pos in routers:
        router = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, router_y, Inches(1.5), Inches(0.7))
        router.fill.solid()
        router.fill.fore_color.rgb = colors['router']
        router.line.color.rgb = RGBColor(50, 80, 150)
        router.line.width = Pt(2)
        add_text_to_shape(router, router_text, Pt(9))
    
    # Link Load Balancer
    llb_y = Inches(2.6)
    llb = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(5.2), llb_y, Inches(3), Inches(0.6))
    llb.fill.solid()
    llb.fill.fore_color.rgb = colors['lb']
    llb.line.width = Pt(2)
    add_text_to_shape(llb, "Link Load Balancer (LLB)", Pt(11), bold=True)
    
    # Perimeter Firewall with IPS
    pf_y = Inches(3.4)
    perimeter_fw = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.5), pf_y, Inches(4.5), Inches(0.7))
    perimeter_fw.fill.solid()
    perimeter_fw.fill.fore_color.rgb = colors['firewall']
    perimeter_fw.line.width = Pt(2)
    add_text_to_shape(perimeter_fw, "Perimeter Firewall - PA-5410 + IPS", Pt(12), bold=True)
    
    # DMZ Zone
    dmz_y = Inches(4.3)
    dmz_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), dmz_y, Inches(3.5), Inches(1.8))
    dmz_zone.fill.solid()
    dmz_zone.fill.fore_color.rgb = colors['dmz']
    dmz_zone.line.color.rgb = RGBColor(200, 50, 50)
    dmz_zone.line.width = Pt(2)
    
    dmz_title = slide.shapes.add_textbox(Inches(0.6), dmz_y + Inches(0.05), Inches(3.3), Inches(0.3))
    dmz_title.text_frame.text = "DMZ Zone (VLAN 103)"
    dmz_title.text_frame.paragraphs[0].font.size = Pt(12)
    dmz_title.text_frame.paragraphs[0].font.bold = True
    dmz_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Sterling File Gateway cluster in DMZ
    sfg1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), dmz_y + Inches(0.5), Inches(1.4), Inches(0.5))
    sfg1.fill.solid()
    sfg1.fill.fore_color.rgb = colors['server']
    add_text_to_shape(sfg1, "SFG-SSP 1", Pt(10))
    
    sfg2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.3), dmz_y + Inches(0.5), Inches(1.4), Inches(0.5))
    sfg2.fill.solid()
    sfg2.fill.fore_color.rgb = colors['server']
    add_text_to_shape(sfg2, "SFG-SSP 2", Pt(10))
    
    sfg_label = slide.shapes.add_textbox(Inches(0.8), dmz_y + Inches(1.1), Inches(2.9), Inches(0.3))
    sfg_label.text_frame.text = "Sterling File Gateway\nSecure File Transfer"
    sfg_label.text_frame.paragraphs[0].font.size = Pt(9)
    sfg_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    sfg_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # WAF
    waf = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(3.2), dmz_y + Inches(1.3), Inches(0.7), Inches(0.4))
    waf.fill.solid()
    waf.fill.fore_color.rgb = RGBColor(255, 150, 150)
    add_text_to_shape(waf, "WAF", Pt(9))
    
    # Core Firewall
    cf_y = Inches(4.6)
    core_fw = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.5), cf_y, Inches(4.5), Inches(0.6))
    core_fw.fill.solid()
    core_fw.fill.fore_color.rgb = colors['firewall']
    core_fw.line.width = Pt(2)
    add_text_to_shape(core_fw, "Core Firewall - FG-400F", Pt(11), bold=True)
    
    # Web/App Zone
    webapp_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.3), dmz_y, Inches(4.5), Inches(1.8))
    webapp_zone.fill.solid()
    webapp_zone.fill.fore_color.rgb = colors['web_app']
    webapp_zone.line.color.rgb = RGBColor(50, 200, 50)
    webapp_zone.line.width = Pt(2)
    
    webapp_title = slide.shapes.add_textbox(Inches(4.4), dmz_y + Inches(0.05), Inches(4.3), Inches(0.3))
    webapp_title.text_frame.text = "Web/Application Zone"
    webapp_title.text_frame.paragraphs[0].font.size = Pt(12)
    webapp_title.text_frame.paragraphs[0].font.bold = True
    webapp_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Server Load Balancer
    slb = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(5.5), dmz_y + Inches(0.4), Inches(2), Inches(0.5))
    slb.fill.solid()
    slb.fill.fore_color.rgb = colors['lb']
    add_text_to_shape(slb, "Server Load Balancer", Pt(10))
    
    # Application components
    app_y = dmz_y + Inches(1.0)
    apps = [
        ("Maximo/MAS", Inches(4.5), Inches(0.9), Inches(0.6)),
        ("Instana", Inches(5.5), Inches(0.9), Inches(0.6)),
        ("Kafka", Inches(6.5), Inches(0.9), Inches(0.6)),
        ("Reporting", Inches(7.5), Inches(0.9), Inches(0.6))
    ]
    
    for app_name, x, w, h in apps:
        app = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, app_y, w, h)
        app.fill.solid()
        app.fill.fore_color.rgb = colors['server']
        add_text_to_shape(app, app_name, Pt(9))
    
    # Management Zone
    mgmt_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), dmz_y, Inches(3.8), Inches(1.8))
    mgmt_zone.fill.solid()
    mgmt_zone.fill.fore_color.rgb = colors['mgmt']
    mgmt_zone.line.color.rgb = RGBColor(200, 200, 50)
    mgmt_zone.line.width = Pt(2)
    
    mgmt_title = slide.shapes.add_textbox(Inches(9.1), dmz_y + Inches(0.05), Inches(3.6), Inches(0.3))
    mgmt_title.text_frame.text = "Management Zone (VLAN 104)"
    mgmt_title.text_frame.paragraphs[0].font.size = Pt(12)
    mgmt_title.text_frame.paragraphs[0].font.bold = True
    mgmt_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Management components
    mgmt_y = dmz_y + Inches(0.5)
    mgmt_apps = [
        ("vPAS NMS", Inches(9.2)),
        ("SevOne", Inches(10.3)),
        ("vDNC", Inches(11.4))
    ]
    
    for app_name, x in mgmt_apps:
        app = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, mgmt_y, Inches(0.9), Inches(0.5))
        app.fill.solid()
        app.fill.fore_color.rgb = colors['server']
        add_text_to_shape(app, app_name, Pt(9))
    
    # Network monitoring details
    mon_details = slide.shapes.add_textbox(Inches(9.2), mgmt_y + Inches(0.7), Inches(3.4), Inches(0.8))
    mon_text = "Network Monitoring:\n• SNMP, SNMP Trap\n• NetFlow/sFlow (UDP 9996)\n• ICMP, WMI"
    mon_details.text_frame.text = mon_text
    mon_details.text_frame.paragraphs[0].font.size = Pt(8)
    set_textbox_color(mon_details)
    
    # MZ/DB Zone
    db_y = Inches(6.3)
    db_zone = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), db_y, Inches(12.3), Inches(0.8))
    db_zone.fill.solid()
    db_zone.fill.fore_color.rgb = colors['mz_db']
    db_zone.line.color.rgb = RGBColor(50, 50, 200)
    db_zone.line.width = Pt(2)
    
    db_components = [
        ("DB2 Database", Inches(1)),
        ("MongoDB", Inches(3)),
        ("OpenShift Platform", Inches(5)),
        ("Storage (16G FC)", Inches(7.5)),
        ("Backup Services", Inches(10))
    ]
    
    for comp_name, x in db_components:
        comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, db_y + Inches(0.15), Inches(1.8), Inches(0.5))
        comp.fill.solid()
        comp.fill.fore_color.rgb = colors['server']
        add_text_to_shape(comp, comp_name, Pt(9))
    
    # Add zone label
    db_label = slide.shapes.add_textbox(Inches(0.6), db_y + Inches(0.02), Inches(4), Inches(0.2))
    db_label.text_frame.text = "MZ/DB Zone (VLAN 105) - Database & Backend Services"
    db_label.text_frame.paragraphs[0].font.size = Pt(10)
    db_label.text_frame.paragraphs[0].font.bold = True
    db_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Security services box
    sec_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.2), Inches(4), Inches(0.3))
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    sec_text = "Security Services: SSO with LDAP | API Key Authentication | Secure Proxy"
    add_text_to_shape(sec_box, sec_text, Pt(9))
    
    # Common services box
    common_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.7), Inches(7.2), Inches(4), Inches(0.3))
    common_box.fill.solid()
    common_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    common_text = "Common Services (VLAN 111): DNS | DHCP | NTP | AD/LDAP"
    add_text_to_shape(common_box, common_text, Pt(9))
    
    # DC/DR indicator
    dc_dr_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.9), Inches(7.2), Inches(3.9), Inches(0.3))
    dc_dr_box.fill.solid()
    dc_dr_box.fill.fore_color.rgb = RGBColor(240, 240, 255)
    dc_dr_text = "Active-Active DC & DR Sites | Identical Architecture"
    add_text_to_shape(dc_dr_box, dc_dr_text, Pt(9))

def create_enhanced_physical_diagram(prs):
    """Create enhanced physical network architecture diagram"""
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "HAL Physical Network Architecture - Spine-Leaf Topology (DC Site)"
    title_frame.paragraphs[0].font.size = Pt(22)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Colors
    colors = {
        'spine': RGBColor(70, 120, 255),       # Dark blue
        'leaf': RGBColor(100, 200, 100),       # Green
        'server': RGBColor(220, 220, 220),     # Gray
        'storage': RGBColor(255, 180, 100),    # Orange
        'router': RGBColor(150, 150, 255),     # Light blue
        'firewall': RGBColor(255, 120, 60),    # Orange-red
        'mgmt': RGBColor(255, 255, 150),      # Light yellow
    }
    
    # Rack representation
    rack_y = Inches(0.7)
    rack_width = Inches(3)
    rack_height = Inches(6.2)
    rack_spacing = Inches(3.2)
    
    rack_names = ["Rack 1 - Network", "Rack 2 - Compute", "Rack 3 - Storage", "Rack 4 - Management"]
    rack_x_positions = [Inches(0.5), Inches(0.5) + rack_spacing, Inches(0.5) + 2*rack_spacing, Inches(0.5) + 3*rack_spacing]
    
    # Draw racks
    for i, (rack_name, x_pos) in enumerate(zip(rack_names, rack_x_positions)):
        # Rack outline
        rack = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos, rack_y, rack_width, rack_height)
        rack.fill.solid()
        rack.fill.fore_color.rgb = RGBColor(245, 245, 245)
        rack.line.color.rgb = RGBColor(100, 100, 100)
        rack.line.width = Pt(2)
        
        # Rack label
        rack_label = slide.shapes.add_textbox(x_pos, rack_y - Inches(0.35), rack_width, Inches(0.3))
        rack_label.text_frame.text = rack_name
        rack_label.text_frame.paragraphs[0].font.size = Pt(12)
        rack_label.text_frame.paragraphs[0].font.bold = True
        rack_label.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        rack_label.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Populate rack contents
        if i == 0:  # Network Rack
            # Spine switches
            for j in range(2):
                spine = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(0.2 + j*0.7), Inches(2.8), Inches(0.6))
                spine.fill.solid()
                spine.fill.fore_color.rgb = colors['spine']
                add_text_to_shape(spine, f"Spine Switch {j+1}\nN9K-C9316D-GX\n32x100G", Pt(9))
            
            # Border leaf switches
            for j in range(2):
                leaf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(1.8 + j*0.6), Inches(2.8), Inches(0.5))
                leaf.fill.solid()
                leaf.fill.fore_color.rgb = colors['leaf']
                add_text_to_shape(leaf, f"Border Leaf {j+1}\nN9K-C9316D-GX", Pt(8))
            
            # Routers
            for j in range(4):
                router_type = "Internet" if j < 2 else "MPLS"
                router = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(3.2 + j*0.5), Inches(2.8), Inches(0.4))
                router.fill.solid()
                router.fill.fore_color.rgb = colors['router']
                add_text_to_shape(router, f"{router_type} Router {(j%2)+1} - ISR 4461", Pt(8))
            
            # Firewalls
            fw_names = ["Perimeter FW - PA-5410", "Core FW - FG-400F"]
            for j, fw_name in enumerate(fw_names):
                fw = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x_pos + Inches(0.1), rack_y + Inches(5.4 + j*0.5), Inches(2.8), Inches(0.4))
                fw.fill.solid()
                fw.fill.fore_color.rgb = colors['firewall']
                add_text_to_shape(fw, fw_name, Pt(8))
                
        elif i == 1:  # Compute Rack
            # Compute leaf switches
            for j in range(2):
                leaf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(0.2 + j*0.6), Inches(2.8), Inches(0.5))
                leaf.fill.solid()
                leaf.fill.fore_color.rgb = colors['leaf']
                add_text_to_shape(leaf, f"Compute Leaf {j+1}\nN9K-C93180YC-EX", Pt(8))
            
            # OpenShift compute nodes
            for j in range(8):
                server = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(1.6 + j*0.55), Inches(2.8), Inches(0.5))
                server.fill.solid()
                server.fill.fore_color.rgb = colors['server']
                server.line.width = Pt(1)
                add_text_to_shape(server, f"OpenShift Node {j+1}\n2x25G NICs, Dual PSU", Pt(8))
                
        elif i == 2:  # Storage Rack
            # Service leaf switches
            for j in range(2):
                leaf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(0.2 + j*0.6), Inches(2.8), Inches(0.5))
                leaf.fill.solid()
                leaf.fill.fore_color.rgb = colors['leaf']
                add_text_to_shape(leaf, f"Service Leaf {j+1}", Pt(8))
            
            # SAN switches
            for j in range(2):
                san_sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(1.6 + j*0.5), Inches(2.8), Inches(0.4))
                san_sw.fill.solid()
                san_sw.fill.fore_color.rgb = RGBColor(200, 150, 255)
                add_text_to_shape(san_sw, f"SAN Switch {j+1} - 16G FC", Pt(8))
            
            # Storage arrays
            storage = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(2.8), Inches(2.8), Inches(2.5))
            storage.fill.solid()
            storage.fill.fore_color.rgb = colors['storage']
            add_text_to_shape(storage, "Storage Array\n16G FC Connectivity\nDual Controllers", Pt(10))
            
            # Backup
            backup = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(5.5), Inches(2.8), Inches(0.5))
            backup.fill.solid()
            backup.fill.fore_color.rgb = RGBColor(180, 180, 180)
            add_text_to_shape(backup, "Backup Storage", Pt(9))
            
        elif i == 3:  # Management Rack
            # DMZ leaf switches
            for j in range(2):
                leaf = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(0.2 + j*0.6), Inches(2.8), Inches(0.5))
                leaf.fill.solid()
                leaf.fill.fore_color.rgb = colors['leaf']
                add_text_to_shape(leaf, f"DMZ Leaf {j+1}", Pt(8))
            
            # Management servers
            mgmt_servers = [
                "vPAS NMS Server",
                "SevOne Server",
                "vDNC Server",
                "Instana Server",
                "Sterling Gateway 1",
                "Sterling Gateway 2",
                "Management Tools"
            ]
            
            for j, server_name in enumerate(mgmt_servers):
                server = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.1), rack_y + Inches(1.6 + j*0.55), Inches(2.8), Inches(0.5))
                server.fill.solid()
                server.fill.fore_color.rgb = colors['mgmt'] if j < 4 else colors['server']
                add_text_to_shape(server, server_name, Pt(8))
    
    # Add connectivity legend
    legend_y = Inches(7.1)
    legend_items = [
        ("100G", colors['spine'], "Spine-Leaf"),
        ("25G", colors['leaf'], "Compute-Leaf"),
        ("16G FC", colors['storage'], "Storage"),
        ("10G", RGBColor(100, 100, 100), "Service"),
        ("1G", RGBColor(150, 150, 150), "Management")
    ]
    
    legend_x = Inches(0.5)
    for speed, color, desc in legend_items:
        # Color box
        color_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x, legend_y, Inches(0.3), Inches(0.2))
        color_box.fill.solid()
        color_box.fill.fore_color.rgb = color
        
        # Text
        text_box = slide.shapes.add_textbox(legend_x + Inches(0.35), legend_y - Inches(0.05), Inches(1.8), Inches(0.3))
        text_box.text_frame.text = f"{speed} - {desc}"
        text_box.text_frame.paragraphs[0].font.size = Pt(9)
        text_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        legend_x += Inches(2.5)
    
    # Add DR site note
    dr_note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), Inches(0.7), Inches(4.8), Inches(0.5))
    dr_note.fill.solid()
    dr_note.fill.fore_color.rgb = RGBColor(240, 240, 255)
    dr_note.line.color.rgb = RGBColor(100, 100, 200)
    add_text_to_shape(dr_note, "DR Site: 3 Racks with Identical Architecture\nActive-Active Configuration", Pt(10))

def create_vlan_ip_slide(prs):
    """Create VLAN and IP scheme slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "VLAN Configuration and Network Segmentation"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Create table
    table_left = Inches(1)
    table_top = Inches(1.2)
    table_width = Inches(11)
    table_height = Inches(5)
    
    # Add table
    table = slide.shapes.add_table(12, 4, table_left, table_top, table_width, table_height).table
    
    # Set column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3)
    
    # Header row
    headers = ["VLAN Name", "VLAN ID", "Purpose", "Key Components"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(50, 50, 150)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    vlan_data = [
        ("Network Management", "101", "Network device management", "Switches, routers, firewalls mgmt interfaces"),
        ("Compute Management", "102", "Server management interfaces", "iLO/iDRAC, IPMI, console access"),
        ("DMZ", "103", "Demilitarized zone", "Sterling File Gateway, WAF"),
        ("Management Zone (MZ)", "104", "Management applications", "vPAS NMS, SevOne, vDNC"),
        ("Database (DB)", "105", "Database services", "DB2, MongoDB, storage access"),
        ("Security Management", "106", "Security device management", "Firewall mgmt, IPS, security tools"),
        ("Internet Link 1", "107", "Primary internet connectivity", "Internet Router 1, ISP link"),
        ("Internet Link 2", "108", "Secondary internet connectivity", "Internet Router 2, ISP link"),
        ("MPLS Link 1", "109", "Primary MPLS connectivity", "MPLS Router 1, WAN link"),
        ("MPLS Link 2", "110", "Secondary MPLS connectivity", "MPLS Router 2, WAN link"),
        ("Common Services", "111", "Shared infrastructure services", "DNS, DHCP, NTP, AD/LDAP")
    ]
    
    for row_idx, (name, vid, purpose, components) in enumerate(vlan_data, 1):
        # Alternate row colors
        if row_idx % 2 == 0:
            for col in range(4):
                table.cell(row_idx, col).fill.solid()
                table.cell(row_idx, col).fill.fore_color.rgb = RGBColor(240, 240, 250)
        
        table.cell(row_idx, 0).text = name
        table.cell(row_idx, 1).text = vid
        table.cell(row_idx, 2).text = purpose
        table.cell(row_idx, 3).text = components
        
        for col in range(4):
            table.cell(row_idx, col).text_frame.paragraphs[0].font.size = Pt(10)
            table.cell(row_idx, col).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Add note
    note_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.5), Inches(11), Inches(0.8))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RGBColor(255, 255, 240)
    note_box.line.color.rgb = RGBColor(200, 200, 100)
    note_text = "Note: IP addressing scheme to be provided by HAL. All VLANs support both IPv4 and IPv6.\nEach VLAN is configured with appropriate ACLs and firewall rules for security."
    add_text_to_shape(note_box, note_text, Pt(10))

def create_security_architecture_slide(prs):
    """Create security architecture slide"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Security Architecture - Defense in Depth"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Security layers
    layer_y = Inches(1.0)
    layer_height = Inches(0.8)
    layer_spacing = Inches(0.9)
    
    security_layers = [
        ("Perimeter Security", "Internet/MPLS edge protection", ["DDoS Protection", "Perimeter Firewall PA-5410", "IPS at network edge"], RGBColor(255, 100, 100)),
        ("DMZ Security", "Secure external-facing services", ["WAF for web applications", "Isolated network segment", "Sterling Gateway hardening"], RGBColor(255, 150, 100)),
        ("Internal Security", "Core network protection", ["Core Firewall FG-400F", "VLAN segmentation", "Inter-zone access control"], RGBColor(255, 200, 100)),
        ("Application Security", "Application-level protection", ["API security", "SSO with LDAP", "Role-based access control"], RGBColor(200, 255, 200)),
        ("Data Security", "Database and storage protection", ["Encryption at rest", "Database access control", "Backup security"], RGBColor(200, 200, 255)),
        ("Management Security", "Infrastructure protection", ["Privileged access management", "Audit logging", "Security monitoring"], RGBColor(255, 200, 255))
    ]
    
    for i, (layer_name, description, features, color) in enumerate(security_layers):
        y_pos = layer_y + (i * layer_spacing)
        
        # Layer box
        layer_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(3), layer_height)
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.width = Pt(2)
        
        # Layer name
        name_box = slide.shapes.add_textbox(Inches(1.1), y_pos + Inches(0.1), Inches(2.8), Inches(0.3))
        name_box.text_frame.text = layer_name
        name_box.text_frame.paragraphs[0].font.size = Pt(12)
        name_box.text_frame.paragraphs[0].font.bold = True
        name_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.1), y_pos + Inches(0.4), Inches(2.8), Inches(0.3))
        desc_box.text_frame.text = description
        desc_box.text_frame.paragraphs[0].font.size = Pt(9)
        desc_box.text_frame.paragraphs[0].font.italic = True
        desc_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        
        # Features
        features_x = Inches(4.5)
        for j, feature in enumerate(features):
            feature_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, features_x + (j * Inches(2.5)), y_pos + Inches(0.1), Inches(2.3), Inches(0.6))
            feature_box.fill.solid()
            feature_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            feature_box.line.color.rgb = color
            add_text_to_shape(feature_box, feature, Pt(9))
    
    # Security compliance box
    compliance_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(6.8), Inches(11), Inches(0.5))
    compliance_box.fill.solid()
    compliance_box.fill.fore_color.rgb = RGBColor(240, 240, 255)
    compliance_text = "Security Compliance: Industry best practices | Regular security audits | Incident response procedures | 24/7 security monitoring"
    add_text_to_shape(compliance_box, compliance_text, Pt(10))

def create_rack_diagrams_slide(prs):
    """Create rack diagrams slide using provided rack-diagram.jpeg"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "HAL Data Center Rack Layout - Physical Infrastructure"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Add the rack diagram image
    try:
        # Insert the rack diagram image
        rack_image = slide.shapes.add_picture('/Users/atyagi/code/hal/rack-diagram.jpeg', 
                                            Inches(1), Inches(0.8), 
                                            Inches(11), Inches(5.5))
    except:
        # Fallback if image cannot be loaded
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(0.8), Inches(11), Inches(5.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(100, 100, 100)
        add_text_to_shape(placeholder, "Rack Diagram Image\n(rack-diagram.jpeg)", Pt(16))
    
    # Add equipment specification table
    table_y = Inches(6.5)
    
    # Equipment specifications
    equipment_specs = [
        ("MAS Pod MH", "Maximo Application Suite on OpenShift", "High Availability Cluster"),
        ("Switch/Core", "Core Network Infrastructure", "Spine-Leaf Architecture"),
        ("Worker", "OpenShift Worker Nodes", "Compute Infrastructure"),
        ("Enterprise Server", "Management & Application Servers", "Sterling, Monitoring, DB"),
        ("Standard Server", "Standard Application Servers", "Supporting Services"),
        ("Sterling Secure Proxy MH", "Sterling File Gateway Cluster", "Secure File Transfer")
    ]
    
    # Create equipment details boxes
    box_width = Inches(5.8)
    box_height = Inches(0.8)
    left_col_x = Inches(0.5)
    right_col_x = Inches(6.8)
    
    for i, (component, description, purpose) in enumerate(equipment_specs):
        if i < 3:  # Left column
            box_x = left_col_x
            box_y = table_y + (i * (box_height + Inches(0.1)))
        else:  # Right column
            box_x = right_col_x
            box_y = table_y + ((i - 3) * (box_height + Inches(0.1)))
        
        # Component box
        comp_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_x, box_y, box_width, box_height)
        comp_box.fill.solid()
        comp_box.fill.fore_color.rgb = RGBColor(245, 245, 255)
        comp_box.line.color.rgb = RGBColor(100, 100, 150)
        comp_box.line.width = Pt(1)
        
        # Component text
        comp_text = f"{component}\n{description}\n{purpose}"
        add_text_to_shape(comp_box, comp_text, Pt(9))
    
    # Add key features box
    features_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    features_box.fill.solid()
    features_box.fill.fore_color.rgb = RGBColor(255, 255, 240)
    features_box.line.color.rgb = RGBColor(200, 200, 100)
    features_text = "Key Features: Redundant Power & Cooling | Hot-Swappable Components | Cable Management | Environmental Monitoring | Remote Access via KVM"
    add_text_to_shape(features_box, features_text, Pt(10))

def add_text_to_shape(shape, text, font_size, bold=False, font_color=None):
    """Add text to a shape with formatting"""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_right = Inches(0.05)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    
    # Default to black font color if not specified
    if font_color is None:
        font_color = RGBColor(0, 0, 0)  # Black
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = font_size
        paragraph.font.bold = bold
        paragraph.font.color.rgb = font_color

def set_textbox_color(textbox, font_color=None):
    """Set font color for textbox content"""
    if font_color is None:
        font_color = RGBColor(0, 0, 0)  # Black
    
    for paragraph in textbox.text_frame.paragraphs:
        paragraph.font.color.rgb = font_color

if __name__ == "__main__":
    create_enhanced_network_presentation()